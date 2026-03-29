import base64
import html
import io
import re
from html.parser import HTMLParser
from typing import Tuple
from urllib.parse import urlparse, unquote

# PDF reader fallback
PdfReader = None
try:
    from pypdf import PdfReader as _PdfReader
    PdfReader = _PdfReader
except Exception:
    try:
        from PyPDF2 import PdfReader as _PdfReader
        PdfReader = _PdfReader
    except Exception:
        PdfReader = None

# Word reader (optional)
_docx = None
try:
    import docx as _docx  # python-docx
except Exception:
    pass

# PowerPoint reader (optional)
_pptx = None
try:
    import pptx as _pptx  # python-pptx
except Exception:
    pass

# Excel reader (optional)
_openpyxl = None
try:
    import openpyxl as _openpyxl
except Exception:
    pass

# Image reader (optional)
_PIL_Image = None
try:
    from PIL import Image as _PIL_Image  # type: ignore
except Exception:
    pass

# EPUB reader (optional)
_ebooklib_epub = None
try:
    from ebooklib import epub as _ebooklib_epub  # type: ignore
except Exception:
    pass


# ---------------------------------------------------------------------------
# Text utilities
# ---------------------------------------------------------------------------

def estimate_tokens(text: str) -> int:
    return max(1, len(text) // 4)


def truncate_to_token_limit(text: str, max_tokens: int) -> Tuple[str, bool]:
    est = estimate_tokens(text)
    if est <= max_tokens:
        return text, False
    ratio = max_tokens / est
    new_len = max(1, int(len(text) * ratio))
    return text[:new_len], True


def normalize_extracted_text(text: str) -> str:
    if not text:
        return ""
    text = text.replace("\x00", "")
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def text_quality_score(text: str) -> int:
    return len(re.sub(r"\s+", "", text or ""))


def truncate_for_print(text: str, max_chars: int) -> str:
    if not text:
        return ""
    if max_chars and len(text) > max_chars:
        return text[:max_chars] + "\n...[truncated]..."
    return text


# ---------------------------------------------------------------------------
# File reading
# ---------------------------------------------------------------------------

def read_file_url(file_url: str) -> bytes:
    parsed = urlparse(file_url)
    if parsed.scheme != "file":
        raise ValueError(f"Not a file URL: {file_url}")
    path = unquote(parsed.path)
    if not path:
        raise ValueError(f"Empty file path in URL: {file_url}")
    with open(path, "rb") as f:
        return f.read()


# ---------------------------------------------------------------------------
# File type detection
# ---------------------------------------------------------------------------

_IMAGE_EXTS = (".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff", ".tif")
_IMAGE_MIME_PREFIXES = ("image/png", "image/jpeg", "image/gif", "image/webp", "image/bmp", "image/tiff")


def detect_file_type(filename: str, content_type: str) -> str:
    """Return one of: 'pdf','word','excel','pptx','markdown','html','image',
    'txt','csv','json','rtf','epub','odt', or '' for unsupported types."""
    name = filename.lower()
    ct = content_type.lower()
    if "pdf" in ct or name.endswith(".pdf"):
        return "pdf"
    if (
        "wordprocessingml" in ct
        or "msword" in ct
        or "macroenabled" in ct
        or name.endswith(".docx")
        or name.endswith(".doc")
        or name.endswith(".docm")
    ):
        return "word"
    if (
        "spreadsheetml" in ct
        or "ms-excel" in ct
        or "excel" in ct
        or name.endswith(".xlsx")
        or name.endswith(".xls")
    ):
        return "excel"
    if (
        "presentationml" in ct
        or name.endswith(".pptx")
        or name.endswith(".ppt")
        or name.endswith(".pptm")
    ):
        return "pptx"
    if name.endswith(".md") or name.endswith(".markdown"):
        return "markdown"
    if (
        "html" in ct
        or name.endswith(".html")
        or name.endswith(".htm")
    ):
        return "html"
    if any(ct.startswith(p) for p in _IMAGE_MIME_PREFIXES) or any(name.endswith(e) for e in _IMAGE_EXTS):
        return "image"
    if "rtf" in ct or name.endswith(".rtf"):
        return "rtf"
    if "epub" in ct or name.endswith(".epub"):
        return "epub"
    if "opendocument.text" in ct or name.endswith(".odt"):
        return "odt"
    if "text/csv" in ct or name.endswith(".csv"):
        return "csv"
    if "application/json" in ct or name.endswith(".json"):
        return "json"
    if "text/plain" in ct or name.endswith(".txt"):
        return "txt"
    return ""


# ---------------------------------------------------------------------------
# PDF extraction
# ---------------------------------------------------------------------------

def extract_pdf_text_pypdf(pdf_bytes: bytes, max_pages: int) -> Tuple[str, int, bool]:
    if PdfReader is None:
        return "", 0, False
    try:
        reader = PdfReader(io.BytesIO(pdf_bytes))
        total_pages = len(reader.pages)
    except Exception:
        return "", 0, False
    use_pages = min(total_pages, max_pages) if max_pages > 0 else total_pages
    chunks = []
    for i in range(use_pages):
        try:
            page_text = reader.pages[i].extract_text() or ""
        except Exception:
            page_text = ""
        if page_text:
            chunks.append(page_text)
    text = normalize_extracted_text("\n\n".join(chunks))
    return text, total_pages, total_pages > use_pages


def extract_pdf_text_pdfminer(pdf_bytes: bytes, max_pages: int) -> str:
    try:
        from pdfminer.high_level import extract_text  # type: ignore
    except Exception:
        return ""
    try:
        kwargs = {}
        if max_pages > 0:
            kwargs["maxpages"] = max_pages
        text = extract_text(io.BytesIO(pdf_bytes), **kwargs) or ""
        return normalize_extracted_text(text)
    except Exception:
        return ""


def extract_pdf_text(pdf_bytes: bytes, max_pages: int) -> Tuple[str, int, bool]:
    pypdf_text, total_pages, truncated_pages = extract_pdf_text_pypdf(pdf_bytes, max_pages)
    pdfminer_text = ""
    if text_quality_score(pypdf_text) < 200:
        pdfminer_text = extract_pdf_text_pdfminer(pdf_bytes, max_pages)
    if text_quality_score(pdfminer_text) > text_quality_score(pypdf_text):
        return pdfminer_text, total_pages, truncated_pages
    if pypdf_text:
        return pypdf_text, total_pages, truncated_pages
    if pdfminer_text:
        return pdfminer_text, total_pages, truncated_pages
    raise RuntimeError(
        "PDF text extraction failed. Install `pypdf` or `pdfminer.six` "
        "with `python3 -m pip install pypdf pdfminer.six`."
    )


# ---------------------------------------------------------------------------
# Word / Excel extraction
# ---------------------------------------------------------------------------

def _find_soffice() -> str:
    """Return path to LibreOffice soffice executable, or ''."""
    import platform
    import shutil
    import os
    system = platform.system()
    if system == "Windows":
        for candidate in [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]:
            if os.path.isfile(candidate):
                return candidate
    elif system == "Darwin":
        path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        if os.path.isfile(path):
            return path
    return shutil.which("soffice") or shutil.which("libreoffice") or ""


def _extract_doc_via_libreoffice(file_bytes: bytes) -> str:
    """Cross-platform fallback: convert .doc to txt via LibreOffice headless."""
    import os
    import subprocess
    import tempfile
    soffice = _find_soffice()
    if not soffice:
        return ""
    with tempfile.TemporaryDirectory() as tmpdir:
        doc_path = os.path.join(tmpdir, "input.doc")
        with open(doc_path, "wb") as f:
            f.write(file_bytes)
        try:
            subprocess.run(
                [soffice, "--headless", "--convert-to", "txt:Text", "--outdir", tmpdir, doc_path],
                capture_output=True,
                timeout=60,
            )
            txt_path = os.path.join(tmpdir, "input.txt")
            if os.path.exists(txt_path):
                with open(txt_path, "rb") as f:
                    return f.read().decode("utf-8", errors="replace")
        except Exception:
            pass
    return ""


def _extract_doc_via_antiword(file_bytes: bytes) -> str:
    """Unix/macOS fallback: extract .doc text via antiword subprocess."""
    import os
    import subprocess
    import tempfile
    tmp = tempfile.NamedTemporaryFile(suffix=".doc", delete=False)
    try:
        tmp.write(file_bytes)
        tmp.close()
        result = subprocess.run(
            ["antiword", tmp.name],
            capture_output=True,
            timeout=30,
        )
        if result.returncode == 0:
            return result.stdout.decode("utf-8", errors="replace")
    except FileNotFoundError:
        pass
    except Exception:
        pass
    finally:
        os.unlink(tmp.name)
    return ""


def extract_word_text(file_bytes: bytes) -> Tuple[str, int, bool]:
    if _docx is None:
        raise RuntimeError(
            "Word text extraction failed. Install `python-docx` "
            "with `python3 -m pip install python-docx`."
        )
    try:
        doc = _docx.Document(io.BytesIO(file_bytes))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = normalize_extracted_text("\n\n".join(paragraphs))
        return text, len(paragraphs), False
    except Exception as e:
        msg = str(e).lower()
        if "zip" not in msg and "not a word file" not in msg:
            raise  # 非格式问题，直接抛出

    # 老式二进制 .doc：依次尝试 LibreOffice（跨平台）→ antiword（Unix/macOS）
    for extractor in (_extract_doc_via_libreoffice, _extract_doc_via_antiword):
        raw = extractor(file_bytes)
        if raw:
            text = normalize_extracted_text(raw)
            para_count = len([p for p in text.split("\n\n") if p.strip()])
            return text, para_count, False

    raise RuntimeError(
        "Old binary .doc format could not be parsed. "
        "Install LibreOffice (https://www.libreoffice.org) or antiword "
        "(`brew install antiword` on macOS / `apt install antiword` on Linux), "
        "or convert the file to .docx first."
    )


def extract_pptx_text(file_bytes: bytes) -> Tuple[str, int, bool]:
    if _pptx is None:
        raise RuntimeError(
            "PowerPoint text extraction failed. Install `python-pptx` "
            "with `python3 -m pip install python-pptx`."
        )
    prs = _pptx.Presentation(io.BytesIO(file_bytes))
    total_slides = len(prs.slides)
    chunks = []
    for slide in prs.slides:
        slide_parts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if line:
                    slide_parts.append(line)
        if slide_parts:
            chunks.append("\n".join(slide_parts))
    text = normalize_extracted_text("\n\n".join(chunks))
    return text, total_slides, False


def _decode_text_bytes(file_bytes: bytes) -> str:
    """Try common encodings in order, fall back to utf-8 with replacement."""
    for enc in ("utf-8-sig", "utf-8", "gbk", "cp1252", "latin-1"):
        try:
            return file_bytes.decode(enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return file_bytes.decode("utf-8", errors="replace")


def extract_txt_text(file_bytes: bytes) -> Tuple[str, int, bool]:
    raw = _decode_text_bytes(file_bytes)
    text = normalize_extracted_text(raw)
    para_count = len([p for p in text.split("\n\n") if p.strip()])
    return text, para_count, False


def extract_markdown_text(file_bytes: bytes) -> Tuple[str, int, bool]:
    raw = _decode_text_bytes(file_bytes)
    text = normalize_extracted_text(raw)
    para_count = len([p for p in text.split("\n\n") if p.strip()])
    return text, para_count, False


def extract_csv_text(file_bytes: bytes, max_rows: int = 500) -> Tuple[str, int, bool]:
    import csv as _csv
    raw = _decode_text_bytes(file_bytes)
    reader = _csv.reader(io.StringIO(raw))
    rows = []
    total_rows = 0
    truncated = False
    for row in reader:
        cells = [c.strip() for c in row if c.strip()]
        if cells:
            rows.append("\t".join(cells))
            total_rows += 1
        if max_rows > 0 and total_rows >= max_rows:
            truncated = True
            break
    text = normalize_extracted_text("\n".join(rows))
    return text, total_rows, truncated


def extract_json_text(file_bytes: bytes) -> Tuple[str, int, bool]:
    import json as _json
    raw = _decode_text_bytes(file_bytes)
    try:
        data = _json.loads(raw)
        text = _json.dumps(data, ensure_ascii=False, indent=2)
    except Exception:
        text = raw
    text = normalize_extracted_text(text)
    return text, 1, False


def extract_rtf_text(file_bytes: bytes) -> Tuple[str, int, bool]:
    try:
        from striprtf.striprtf import rtf_to_text  # type: ignore
    except Exception:
        raise RuntimeError(
            "RTF text extraction failed. Install `striprtf` "
            "with `python3 -m pip install striprtf`."
        )
    raw = _decode_text_bytes(file_bytes)
    text = normalize_extracted_text(rtf_to_text(raw))
    para_count = len([p for p in text.split("\n\n") if p.strip()])
    return text, para_count, False


def extract_epub_text(file_bytes: bytes) -> Tuple[str, int, bool]:
    import os
    import tempfile
    if _ebooklib_epub is None:
        raise RuntimeError(
            "EPUB text extraction failed. Install `ebooklib` "
            "with `python3 -m pip install ebooklib`."
        )
    from ebooklib import ITEM_DOCUMENT  # type: ignore
    tmp = tempfile.NamedTemporaryFile(suffix=".epub", delete=False)
    try:
        tmp.write(file_bytes)
        tmp.close()
        book = _ebooklib_epub.read_epub(tmp.name)
        chunks = []
        for item in book.get_items_of_type(ITEM_DOCUMENT):
            part, _, _ = extract_html_text(item.get_content())
            if part.strip():
                chunks.append(part)
        text = normalize_extracted_text("\n\n".join(chunks))
    except Exception:
        tmp.close()
        raise
    finally:
        try:
            os.unlink(tmp.name)
        except FileNotFoundError:
            pass
    return text, len(chunks), False


def extract_odt_text(file_bytes: bytes) -> Tuple[str, int, bool]:
    try:
        from odf.opendocument import load as odf_load  # type: ignore
        from odf.text import P  # type: ignore
    except Exception:
        raise RuntimeError(
            "ODT text extraction failed. Install `odfpy` "
            "with `python3 -m pip install odfpy`."
        )
    doc = odf_load(io.BytesIO(file_bytes))
    paragraphs = []
    for para in doc.getElementsByType(P):
        parts = []
        for node in para.childNodes:
            if hasattr(node, "data"):
                parts.append(node.data)
            elif hasattr(node, "childNodes"):
                for child in node.childNodes:
                    if hasattr(child, "data"):
                        parts.append(child.data)
        line = "".join(parts).strip()
        if line:
            paragraphs.append(line)
    text = normalize_extracted_text("\n\n".join(paragraphs))
    return text, len(paragraphs), False


class _TextExtractor(HTMLParser):
    """Strip HTML tags and decode entities, preserving block-level whitespace."""
    _BLOCK_TAGS = {
        "p", "div", "br", "li", "tr", "td", "th", "h1", "h2", "h3",
        "h4", "h5", "h6", "blockquote", "pre", "article", "section",
        "header", "footer", "nav", "aside", "main",
    }
    _SKIP_TAGS = {"script", "style", "noscript", "head"}

    def __init__(self) -> None:
        super().__init__()
        self._parts: list[str] = []
        self._skip_depth = 0

    def handle_starttag(self, tag: str, attrs: list) -> None:
        if tag in self._SKIP_TAGS:
            self._skip_depth += 1
        elif tag in self._BLOCK_TAGS and self._parts and self._parts[-1] != "\n":
            self._parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in self._SKIP_TAGS:
            self._skip_depth = max(0, self._skip_depth - 1)
        elif tag in self._BLOCK_TAGS:
            self._parts.append("\n")

    def handle_data(self, data: str) -> None:
        if self._skip_depth == 0:
            self._parts.append(data)

    def handle_entityref(self, name: str) -> None:
        if self._skip_depth == 0:
            self._parts.append(html.unescape(f"&{name};"))

    def handle_charref(self, name: str) -> None:
        if self._skip_depth == 0:
            self._parts.append(html.unescape(f"&#{name};"))

    def get_text(self) -> str:
        return "".join(self._parts)


def extract_html_text(file_bytes: bytes, encoding: str = "") -> Tuple[str, int, bool]:
    """Extract plain text from HTML bytes. Returns (text, paragraph_count, False)."""
    if not encoding:
        # Try to detect encoding from BOM or meta charset
        raw = file_bytes[:4096]
        if raw.startswith(b"\xef\xbb\xbf"):
            encoding = "utf-8-sig"
        else:
            m = re.search(rb'charset=["\']?([\w-]+)', raw, re.IGNORECASE)
            encoding = m.group(1).decode("ascii", errors="ignore") if m else "utf-8"
    try:
        markup = file_bytes.decode(encoding, errors="replace")
    except (LookupError, UnicodeDecodeError):
        markup = file_bytes.decode("utf-8", errors="replace")

    parser = _TextExtractor()
    parser.feed(markup)
    raw_text = parser.get_text()
    text = normalize_extracted_text(raw_text)
    para_count = len([p for p in text.split("\n\n") if p.strip()])
    return text, para_count, False


def extract_excel_text(file_bytes: bytes, max_rows: int = 500) -> Tuple[str, int, bool]:
    if _openpyxl is None:
        raise RuntimeError(
            "Excel text extraction failed. Install `openpyxl` "
            "with `python3 -m pip install openpyxl`."
        )
    wb = _openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    chunks = []
    total_rows = 0
    truncated = False
    for sheet in wb.worksheets:
        sheet_rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                sheet_rows.append("\t".join(cells))
                total_rows += 1
            if max_rows > 0 and total_rows >= max_rows:
                truncated = True
                break
        if sheet_rows:
            chunks.append(f"[Sheet: {sheet.title}]\n" + "\n".join(sheet_rows))
        if truncated:
            break
    wb.close()
    text = normalize_extracted_text("\n\n".join(chunks))
    return text, total_rows, truncated


# ---------------------------------------------------------------------------
# Image resizing & encoding
# ---------------------------------------------------------------------------

_SAVE_FMT = {"JPEG": ("JPEG", "image/jpeg"), "PNG": ("PNG", "image/png"),
             "GIF": ("PNG", "image/png"), "WEBP": ("WEBP", "image/webp"),
             "BMP": ("PNG", "image/png"), "TIFF": ("PNG", "image/png"),
             "TIF": ("PNG", "image/png")}


def resize_and_encode_image(file_bytes: bytes, max_long_side: int = 1280) -> Tuple[str, str]:
    """缩放图片到 max_long_side 以内，返回 (base64字符串, mime类型)。
    需要 Pillow：pip install Pillow
    """
    if _PIL_Image is None:
        raise RuntimeError(
            "图片处理需要 Pillow，请执行：python3 -m pip install Pillow"
        )
    img = _PIL_Image.open(io.BytesIO(file_bytes))
    orig_fmt = (img.format or "JPEG").upper()
    save_fmt, mime = _SAVE_FMT.get(orig_fmt, ("PNG", "image/png"))

    w, h = img.size
    if max(w, h) > max_long_side:
        ratio = max_long_side / max(w, h)
        img = img.resize((max(1, int(w * ratio)), max(1, int(h * ratio))),
                         _PIL_Image.LANCZOS)

    # JPEG 不支持 RGBA/P 模式
    if save_fmt == "JPEG" and img.mode not in ("RGB", "L"):
        img = img.convert("RGB")

    buf = io.BytesIO()
    if save_fmt == "JPEG":
        img.save(buf, format="JPEG", quality=85, optimize=True)
    else:
        img.save(buf, format=save_fmt)

    return base64.b64encode(buf.getvalue()).decode("ascii"), mime
